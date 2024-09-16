from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Crear una presentación de PowerPoint basada en el estilo observado
presentation = Presentation()

# Definir el estilo de fuente y color basado en el documento adjunto
font_name = 'Arial'
title_font_size = Pt(44)
content_font_size = Pt(24)
title_color = RGBColor(0, 51, 102)  # Azul oscuro
content_color = RGBColor(0, 0, 0)  # Negro

# Títulos y contenidos para las diapositivas
slide_info = [
    ("Introducción a Machine Learning con Python",
     "Machine Learning es una rama de la inteligencia artificial que se centra en desarrollar algoritmos que aprenden de los datos para hacer predicciones o tomar decisiones automáticas."),
    ("Tipos de Machine Learning",
     "1. Aprendizaje Supervisado: Utiliza datos etiquetados para entrenar modelos que predicen o clasifican datos no vistos.\n"
     "Ejemplos: Regresión, Clasificación.\n\n"
     "2. Aprendizaje No Supervisado: No utiliza etiquetas y busca patrones o estructuras ocultas en los datos.\n"
     "Ejemplos: Clustering, Reducción de Dimensionalidad."),
    ("Aprendizaje Supervisado",
     "El aprendizaje supervisado utiliza un conjunto de datos de entrenamiento con entradas y salidas conocidas para aprender a mapear entradas a salidas.\n\n"
     "Ejemplos:\n"
     "- Regresión Lineal: Predice un valor continuo.\n"
     "- Clasificación: Asigna una clase a cada entrada."),
    ("Aprendizaje No Supervisado",
     "El aprendizaje no supervisado utiliza datos sin etiquetas para identificar patrones o relaciones subyacentes.\n\n"
     "Ejemplos:\n"
     "- Clustering: Agrupa datos en clusters basados en similitudes.\n"
     "- Reducción de Dimensionalidad: Simplifica los datos manteniendo solo las características más importantes."),
    ("Regresión Lineal con Scikit-Learn",
     "La Regresión Lineal modela la relación entre una variable dependiente y una o más variables independientes.\n\n"
     "Fórmula Matemática:\n"
     "y = β₀ + β₁x₁ + β₂x₂ + ... + βₙxₙ + ε\n"
     "donde:\n"
     "y = variable dependiente,\n"
     "x = variables independientes,\n"
     "β = coeficientes,\n"
     "ε = error."),
    ("Regresión Logística con Scikit-Learn",
     "La Regresión Logística se utiliza para problemas de clasificación binaria.\n\n"
     "Fórmula Matemática:\n"
     "p = 1 / (1 + e^(-z)), donde z = β₀ + β₁x₁ + β₂x₂ + ... + βₙxₙ\n"
     "donde:\n"
     "p = probabilidad de la clase positiva,\n"
     "x = variables independientes,\n"
     "β = coeficientes."),
]

# Crear las diapositivas con los títulos y contenidos
for title, content in slide_info:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    # Establecer estilo de título
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = title_font_size
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.name = font_name
    title_box.text_frame.paragraphs[0].font.color.rgb = title_color

    # Establecer estilo de contenido
    content_box.text = content
    for paragraph in content_box.text_frame.paragraphs:
        paragraph.font.size = content_font_size
        paragraph.font.name = font_name
        paragraph.font.color.rgb = content_color

# Guardar la presentación en un archivo
presentation_path = "Machine_Learning_with_Python_Course_Styled.pptx"
presentation.save(presentation_path)

presentation_path